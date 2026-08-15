"""HACKATHON_PROPOSAL.mdの内容から都知事杯提出用スライド(HACKATHON_SLIDES.pptx)を生成する。

内容はHACKATHON_PROPOSAL.mdと同期させること(数値・技術記述を修正したら、
このスクリプトの該当スライドも合わせて直してから再生成する)。
文字量が多い章(3・6・8・9章)は判読性のため前後半2枚に分割している。

実行: uv run --with python-pptx python build_hackathon_slides.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PRIMARY = RGBColor(0x1B, 0x5E, 0x7A)
PRIMARY_DARK = RGBColor(0x12, 0x40, 0x54)
ACCENT = RGBColor(0xE8, 0x5D, 0x3A)
BG_LIGHT = RGBColor(0xF7, 0xF9, 0xFB)
TEXT_DARK = RGBColor(0x2A, 0x2A, 0x2A)
TEXT_GRAY = RGBColor(0x8A, 0x93, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_HEAD = "Yu Gothic UI"
FONT_BODY = "Yu Gothic"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def set_run(p, text, size, color, bold=False, font=FONT_BODY, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return r


BULLET = "●"
SUBBULLET = "‐"


def add_bullets(tf, items, size, line_spacing=1.18, space_after=10,
                 color=TEXT_DARK, accent_color=ACCENT):
    """items: list of (level:int, text:str, numbered:str|None)."""
    first = True
    for level, text, numbered in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        indent = level * 0.35
        p.level = 0
        mark = numbered if numbered else (BULLET if level == 0 else SUBBULLET)
        set_run(p, f"{'　' * (level * 2)}{mark} ", size, accent_color if level == 0 else TEXT_GRAY,
                bold=(level == 0), font=FONT_BODY)
        set_run(p, text, size, color, font=FONT_BODY)


def add_header(slide, chapter_no, title):
    fill_bg(slide, BG_LIGHT)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    bar.shadow.inherit = False

    if chapter_no:
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.45), Inches(0.2), Inches(0.65), Inches(0.65))
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT
        badge.line.fill.background()
        badge.shadow.inherit = False
        btf = badge.text_frame
        btf.word_wrap = False
        btf.margin_left = 0
        btf.margin_right = 0
        btf.margin_top = 0
        btf.margin_bottom = 0
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        set_run(bp, chapter_no, 20, WHITE, bold=True, font=FONT_HEAD)
        title_left = Inches(1.3)
    else:
        title_left = Inches(0.5)

    _, ttf = add_textbox(slide, title_left, Inches(0.18), Inches(11.5), Inches(0.7), anchor=MSO_ANCHOR.MIDDLE)
    tp = ttf.paragraphs[0]
    set_run(tp, title, 26, WHITE, bold=True, font=FONT_HEAD)


def add_footer(slide, page_no):
    _, tf = add_textbox(slide, Inches(0.5), Inches(7.12), Inches(6), Inches(0.3))
    p = tf.paragraphs[0]
    set_run(p, "あんしんつうがくナビ ／ AIあんぜんせんせい", 10, TEXT_GRAY, font=FONT_BODY)
    _, tf2 = add_textbox(slide, Inches(12.3), Inches(7.12), Inches(0.6), Inches(0.3))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    set_run(p2, str(page_no), 10, TEXT_GRAY, font=FONT_BODY)


def content_slide(prs, page_no, chapter_no, title, items, font_size=18):
    slide = blank_slide(prs)
    add_header(slide, chapter_no, title)
    _, tf = add_textbox(slide, Inches(0.7), Inches(1.35), Inches(11.9), Inches(5.6))
    add_bullets(tf, items, font_size)
    add_footer(slide, page_no)
    return slide


def title_slide(prs):
    slide = blank_slide(prs)
    fill_bg(slide, PRIMARY)
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.55), Inches(2.2), Inches(0.08))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT
    accent_bar.line.fill.background()
    accent_bar.shadow.inherit = False
    accent_bar.left = Inches(5.57)

    _, tf = add_textbox(slide, Inches(1), Inches(2.5), Inches(11.33), Inches(1.4), anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p, "あんしんつうがくナビ ／ AIあんぜんせんせい", 40, WHITE, bold=True, font=FONT_HEAD)

    _, tf2 = add_textbox(slide, Inches(1), Inches(3.85), Inches(11.33), Inches(0.7), anchor=MSO_ANCHOR.MIDDLE)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    set_run(p2, "事故データとAIで、通学路のキケンを親子で見える化", 20, ACCENT, italic=True, font=FONT_BODY)

    _, tf3 = add_textbox(slide, Inches(1), Inches(6.4), Inches(11.33), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    set_run(p3, "都知事杯オープンデータ・ハッカソン2026 提出資料", 14, RGBColor(0xC9, 0xDA, 0xE2), font=FONT_BODY)
    return slide


def closing_slide(prs, page_no):
    slide = blank_slide(prs)
    fill_bg(slide, PRIMARY)
    _, tf = add_textbox(slide, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.6))
    p = tf.paragraphs[0]
    set_run(p, "12. まとめ(一言で)", 22, ACCENT, bold=True, font=FONT_HEAD)

    _, tf2 = add_textbox(slide, Inches(1.3), Inches(2.5), Inches(10.7), Inches(3), anchor=MSO_ANCHOR.MIDDLE)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    p2.line_spacing = 1.35
    set_run(
        p2,
        "「オープンデータ×AIをもとに、データ作成、追加学習を行い、"
        "“なんとなく危ない”を“具体的にここが危ない、こう気をつけよう”に変える、"
        "親子のための通学路安全サービス」",
        24, WHITE, bold=True, font=FONT_HEAD,
    )
    add_footer(slide, page_no)
    return slide


def build():
    prs = new_prs()
    page = 1

    title_slide(prs)
    page += 1

    content_slide(prs, page, "1", "タイトル・キャッチコピー", [
        (0, "プロジェクト名: あんしんつうがくナビ ／ AIあんぜんせんせい", None),
        (0, "キャッチコピー案", None),
        (1, "事故データとAIで、通学路のキケンを親子で見える化", None),
        (1, "「なんとなく危ない」を、データで「ここが危ない理由」に変える", None),
    ], font_size=22)
    page += 1

    content_slide(prs, page, "2", "課題背景(社会課題の提示)", [
        (0, "毎年、通学路での歩行者事故が後を絶たない。保護者や学校は「なんとなく」の感覚でしか通学路の安全性を判断できていない", None),
        (0, "ハザードマップ的な静的な注意喚起はあるが、自宅〜学校の実際のルート単位でリスクを定量的に示すツールは少ない", None),
        (0, "警察庁の交通事故統計オープンデータは公開されているが、生データのままでは一般の保護者・子どもが読み解けない", None),
        (0, "まだ事故が起きていないだけの「隠れた危険地帯」は事故統計だけでは見つけられない", None),
        (0, "→ オープンデータを「誰でも使える意思決定ツール」に翻訳し、実データが少ない場所の潜在リスクも補う", None),
    ], font_size=19)
    page += 1

    content_slide(prs, page, "3", "解決策の概要(1/2)", [
        (0, "自宅・学校の位置をマップ上でピン留めするだけで、通学路沿いの過去の事故を自動集計・地図表示", "①"),
        (0, "地理院タイルの地図画像から「事故が起きやすい地形・道路パターン」をCNNが学習し、実データがまだ少ない場所も含めて危険度を推定(件数のPoisson回帰、単なる◯✕分類ではない)", "②"),
    ], font_size=21)
    page += 1

    content_slide(prs, page, "3", "解決策の概要(2/2)", [
        (0, "Claude(LLM)が危険度モデル・事故データ・信号や横断歩道などの実地点情報をツールとして呼び出し、保護者と子どもにも分かるやさしい言葉で「どこが危ないか」「なぜ危ないか」「どう気をつければよいか」を対話形式で説明", "③"),
        (0, "実際の事故地点・AI予測ヒートマップ・危険な交差点や狭い道路の検出結果を1枚の地図に重ねて可視化", "④"),
    ], font_size=21)
    page += 1

    content_slide(prs, page, "4", "活用するオープンデータ", [
        (0, "警察庁 交通事故統計情報オープンデータ: 東京都内の歩行者事故、2019年〜2024年(約6年分)", None),
        (0, "国土地理院 地理院タイル(std/pale): 危険度モデルの学習用地図画像", None),
        (0, "国土数値情報(国土交通省): 行政区域(N03、都境界の判定)、鉄道駅(N02、最寄り駅距離の算出)", None),
        (0, "PLATEAU(3D都市モデル、G空間情報センター): 23区内の交差点の見通しを実際の3D建物データから計算し、AIの説明の根拠として活用", None),
        (0, "道路種別・信号機・横断歩道等の位置情報を集計統計として活用", None),
    ], font_size=19)
    page += 1

    content_slide(prs, page, "5", "主な機能・デモの流れ", [
        (0, "マップを開く(東京都内、事故データがクラスタ表示された状態)", "1."),
        (0, "「🏡 おうちをきめる」→「🏫 がっこうをきめる」でピンを設置", "2."),
        (0, "ルートが自動描画され、ルート沿いの事故件数が即座に表示される", "3."),
        (0, "「🤖 あんしんせんせいにきいてみる」を押す — Claudeが複数ツールを呼び出しやさしい日本語で解説を生成、地図にAI予測ヒートマップ・実際の事故地点・信号のない横断歩道・狭い道路区間を重ねて表示", "4."),
        (0, "「実際の事故が多い場所」と「AIが危ないと予測した場所」が一致する/しないケースを見せ、両方の情報を持つ価値を伝える", "5."),
    ], font_size=18)
    page += 1

    content_slide(prs, page, "6", "技術構成(アーキテクチャ)(1/2)", [
        (0, "フロントエンド: HTML/CSS/Vanilla JS、Leaflet.js(地図描画・クラスタ・ルーティング)。子ども・保護者向けにひらがな中心のやさしいUI", None),
        (0, "危険度推定モデル(ml_risk_model/)", None),
        (1, "地理院タイル画像(pale版)を500mメッシュで取得し、MobileNetV2ベースの転移学習でPoisson回帰ヘッド(件数予測)を構築", None),
        (1, "空間ブロック分割でtrain/evalを分離(隣接セルの情報漏洩を防止)、3シードアンサンブルで安定化", None),
        (1, "車道延長(道路網の集計統計)を曝露量(オフセット)として組み込み", None),
    ], font_size=18)
    page += 1

    content_slide(prs, page, "6", "技術構成(アーキテクチャ)(2/2)", [
        (0, "AI対話バックエンド: Claude API(Tool Use)が「危険度スコア取得」「近隣事故データ取得」「地点の信号・横断歩道・視界情報取得」をツールとして呼び出し、実データに基づいた説明文を生成", None),
        (0, "インフラ: Cloudflare Workers + Containers(本ハッカソンのCloudflare特典を活用)。フロントは静的配信、推論・LLM連携部分はコンテナ上で実行", None),
        (0, "データパイプライン: 警察庁CSV→GeoJSON変換、地理院タイル取得→CNN学習データ化、道路網データ→集計統計への変換、PLATEAU→視界特徴計算まで、全工程を自前スクリプトで構築", None),
        (0, "プライバシー: 自宅・学校の位置情報はリクエストのたびにその場で処理するのみで、サーバー側のデータベースやログに保存しない設計", None),
    ], font_size=18)
    page += 1

    content_slide(prs, page, "7", "独自性・新規性(差別化ポイント)", [
        (0, "既製のAI APIにリスク判定を丸投げするのではなく、転移学習をベースに自前でデータ収集・特徴量設計・ファインチューニング・評価プロトコルまで一気通貫で構築(単にAI APIを呼び出すだけのアプリではない)", "(1)"),
        (0, "単なる事故マップではなく、AIによる危険度“予測”と実データを重ねて比較できる。まだ事故が記録されていない場所の潜在リスクも地図画像パターンの学習によって補完", "(2)"),
        (0, "判断根拠を「事実」で示す設計。信号の有無・横断歩道の種類・道路の見通し(PLATEAU実測値)など、誰でも現地で確認できる実データに基づいて生成する", "(3)"),
        (0, "評価に学術的な誠実さ: 事前登録方式・複数シードでの再現性確認・bootstrap法による統計的検証", "(4)"),
    ], font_size=17)
    page += 1

    content_slide(prs, page, "8", "学習データの作成規模(1/2)", [
        (0, "事故データ: 警察庁オープンデータから東京都内の歩行者事故を抽出、2019年〜2024年の約6年分・29,284件", None),
        (0, "学習・評価用地図画像: 500mメッシュで空間ブロック分割(隣接セルの情報漏洩を防ぐ設計)し、学習3,803セル・評価896セルを地理院タイルから1枚ずつ生成。学習セットはデータ拡張(回転・反転・明暗ジッター)で15,212枚まで拡張", None),
    ], font_size=20)
    page += 1

    content_slide(prs, page, "8", "学習データの作成規模(2/2)", [
        (0, "位置情報の整備: 国土数値情報から都内駅データ・行政区域データを整備し、最寄り駅距離の算出・都境界の判定に利用", None),
        (0, "見通し(視界)データ: PLATEAU(3D都市モデル)から23区内1,917セル分の交差点見通しを実際の建物形状データから算出", None),
        (0, "地図画像の取得はレート制限(0.5秒/枚)を守った逐次リクエストで行っており、学習用データセット一式の生成だけで数時間規模の作業時間を要した", None),
    ], font_size=20)
    page += 1

    content_slide(prs, page, "9", "開発における試行錯誤(1/2)", [
        (0, "目視評価と実際の学習結果が逆転: 人間の目に情報量が多く見える地図スタイルの方が、CNNの学習には必ずしも有利ではないと実験で判明。複数の候補スタイルで実際に学習・比較し、精度だけでなく複数シード間のばらつき(安定性)も根拠に最終採用スタイルを選び直した", None),
        (0, "見出し数値の水増しに自ら気づき下方修正: 道路がほとんどない場所を含めた全域評価の数値(残差Spearman0.439)は底上げされていたと判明。都市部に限定して再評価し、より正直な数字(0.28、ペアワイズ一致率71%)を公式の見出し数値として採用した", None),
    ], font_size=17)
    page += 1

    content_slide(prs, page, "9", "開発における試行錯誤(2/2)", [
        (0, "特徴量計算バグの自己発見: 評価用ベースラインの入力(建造物密度の代理指標)が、緑地・水域以外を一律「建造物」に誤分類するバグを目視確認で発見。修正後に評価を再実行し、影響範囲を数値で確認してから最終数値を確定させた", None),
        (0, "判定基準の後出し変更をしない運用: 「AIモデルが単純なベースラインを上回るか」の判定基準は結果を見る前に固定し、途中で緩めたり変えたりしない", None),
        (0, "一度出した数値を鵜呑みにせず、追加の目視確認・独立な角度からの再検証を重ねたプロセスの再現性・誠実さも本提案の強み", None),
    ], font_size=17)
    page += 1

    content_slide(prs, page, "10", "今後の展望・拡張性", [
        (0, "東京都オープンデータ(学校位置・道路交通センサス等)を組み合わせた特徴量追加によるさらなる精度向上", None),
        (0, "自治体・学校向けの集計レポート機能(通学路点検・PTA活動での活用)", None),
        (0, "地域住民からの「ヒヤリハット」投稿機能によるデータ拡充", None),
        (0, "23区外(多摩地域)への視界特徴(PLATEAU)適用範囲拡大", None),
    ], font_size=20)
    page += 1

    content_slide(prs, page, "11", "チーム体制・役割分担", [
        (0, "個人参加", None),
        (0, "企画・データ収集/前処理・機械学習モデルの設計と学習・評価設計・バックエンド・フロントエンド・インフラ(Cloudflare)デプロイまで、一人で一気通貫で開発した", None),
    ], font_size=22)
    page += 1

    closing_slide(prs, page)

    prs.save("HACKATHON_SLIDES.pptx")
    print(f"generated {len(prs.slides)} slides -> HACKATHON_SLIDES.pptx")


if __name__ == "__main__":
    build()
